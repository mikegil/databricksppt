from pathlib import Path
from os import path
from enum import Enum
import numbers
from collections.abc import Iterable
import io
import base64

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt
from itertools import islice
import pandas as pd
import numpy as np


class CHART_TYPE(Enum):
    AREA = 'Area'
    AREA_STACKED = 'Area-Stacked'
    AREA_STACKED_100 = 'Area-Stacked-100'
    BAR = 'Bar'
    BAR_STACKED = 'Bar-Stacked'
    BAR_STACKED_100 = 'Bar-Stacked-100'
    COLUMN = 'Column'
    COLUMN_STACKED = 'Column-Stacked'
    COLUMN_STACKED_100 = 'Column-Stacked-100'
    LINE = 'Line'
    LINE_STACKED = 'Line-Stacked'
    LINE_STACKED_100 = 'Line-Stacked-100'
    LINE_MARKED = 'Line-Marked'
    LINE_MARKED_STACKED = 'Line-Marked-Stacked'
    LINE_MARKED_STACKED_100 = 'Line-Marked-Stacked-100'
    DOUGHNUT = 'Doughnut'
    DOUGHNUT_EXPLODED = 'Doughnut-Exploded'
    PIE = 'Pie'
    PIE_EXPLODED = 'Pie-Exploded'
    RADAR = 'Radar'
    RADAR_FILLED = 'Radar-Filled'
    RADAR_MARKED = 'Radar-Marked'
    XY_SCATTER = 'XY-Scatter'
    XY_SCATTER_LINES = 'XY-Scatter-Lines'
    XY_SCATTER_LINES_SMOOTHED = 'XY-Scatter-Lines-Smoothed'
    XY_SCATTER_LINES_MARKED = 'XY-Scatter-Lines-Marked'
    XY_SCATTER_LINES_MARKED_SMOOTHED = 'XY-Scatter-Lines-Marked-Smoothed'
    BUBBLE = 'Bubble'
    TABLE = 'Table'


class LEGEND_POSITION(Enum):
    BOTTOM = 'Bottom'
    CORNER = 'Corner'
    LEFT = 'Left'
    NONE = 'None'
    RIGHT = 'Right'
    TOP = 'Top'


def toPPT(presentation):
    ppt = __create_presentation(presentation)
    if ppt is None or isinstance(ppt, str):
        return 'Could\'t create PPT'

    slide_count = 0
    body_font = presentation.get('body_font')
    if body_font is None:
        body_font = dict(
            name='Verdana',
            size=10
        )
    for slide in presentation.get('slides'):
        if (slide.get('body_font') is None):
            slide['body_font'] = body_font
        slide_body_font = slide.get('body_font')
        slide_count += 1
        new_slide = __create_slide(ppt, slide)
        if new_slide is None or isinstance(new_slide, str):
            return 'Failed to create slide {}: {}'.format(slide_count, new_slide)

        chart_count = 0
        for chart in slide.get('charts'):
            if (chart.get('body_font') is None):
                chart['body_font'] = body_font
            chart_count += 1
            placeholder_num = chart.get('placeholder_num')
            if placeholder_num is not None and placeholder_num > 0:
                placeholder = __get_placeholder(new_slide, placeholder_num)
            else:
                chart_num = slide.get('chart_num', 1)
                placeholder = __get_chart(new_slide, chart_num)

            if placeholder is None or isinstance(placeholder, str):
                return 'Failed to create placeholder for chart {} in slide {}: {}'.format(chart_count, slide_count, placeholder)

            new_chart = __insert_object(new_slide, placeholder, chart)
            if isinstance(new_chart, str):
                return 'Failed to create chart {} in slide {}: {}'.format(chart_count, slide_count, new_chart)

    return ppt


def toBase64URL(pres):
    # Create string shell to insert the base64-encoded data
    output_str = "<a href='data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{}'>Download here</a>"
    # Create a new byte stream to save to
    stream = io.BytesIO()
    # Save the presentation content to the byte stream
    pres.save(stream)
    # Base64 encode the stream and convert to base64 ascii
    encoded = base64.b64encode(stream.getvalue()).decode()

    return output_str.format(encoded)


def __create_presentation(slideInfo):
    template = slideInfo.get('template')
    if (template is not None):
        if (not isinstance(template, str)):
            template = None
        else:
            if (not path.isfile(template)):
                template = None

    return Presentation(template)


def __create_slide(ppt, slide):
    slide_num = slide.get('slide_num', 0)
    layout_num = slide.get('layout_num', 1)
    title = slide.get('title')

    if (len(ppt.slide_layouts) <= layout_num):
        return 'Layout number {} is outside the number of layouts found in this PPT [{}]'.format(layout_num, len(ppt.slide_layouts))

    if slide_num == 0:
        new_slide = ppt.slides.add_slide(ppt.slide_layouts[layout_num])
    else:
        if len(ppt.slides) >= slide_num:
            new_slide = ppt.slides[slide_num-1]
        else:
            return 'Slide number {} is outside the number of slides found in this PPT [{}]'.format(slide_num, len(ppt.slides))

    if new_slide.shapes.title is not None:
        new_slide.shapes.title.text = title

    return new_slide


def __get_placeholder(slide, placeholder_num):
    if len(slide.placeholders) < placeholder_num or placeholder_num <= 0:
        return 'Placeholder number {} outside the number of placeholders found in this slide [{}]'.format(placeholder_num, len(slide.placeholders))

    placeholderIdx = []

    for shape in slide.placeholders:
        placeholderIdx.append(shape.placeholder_format.idx)

    placeholder = slide.placeholders[placeholderIdx[placeholder_num-1]]

    # Remove empty placeholder
    sp = placeholder._sp
    sp.getparent().remove(sp)

    return placeholder


def __get_chart(slide, chart_num):
    if chart_num == 0:
        return 'Neither placeholder_number, nor chart_number were specified for this slide'

    charts_found = 0

    for shape in slide.shapes:
        if shape.has_chart:
            charts_found += 1
        if charts_found == chart_num:
            shape.element.getparent().remove(shape.element)
            return shape

    return 'Chart number {} is outside the number of charts found in this slide [{}]'.format(chart_num, charts_found)


def __infer_category_labels(data):
    for dataframe in data:
        firstCol = dataframe.iloc[:, 0]
        for cell in firstCol:
            if not isinstance(cell, numbers.Number):
                return True

    return False


def __infer_series_labels(data):
    for dataframe in data:
        for col in dataframe.columns:
            if not isinstance(col, numbers.Number):
                return True

    return False


def __transpose_data(chartInfo):
    transposed_data = []

    for dataframe in chartInfo['data']:
        if not isinstance(dataframe, pd.DataFrame):
            return chartInfo

        if chartInfo['first_column_as_labels'] and chartInfo['column_names_as_labels']:
            indexColName = dataframe.columns[0]
            df = dataframe.set_index(
                dataframe.columns[0]).transpose().reset_index()
            df.rename(columns={'index': indexColName}, inplace=True)
        elif chartInfo['column_names_as_labels']:
            df = dataframe.transpose().reset_index()
        elif chartInfo['first_column_as_labels']:
            df = dataframe.set_index(dataframe.columns[0]).transpose()
        else:
            df = dataframe.transpose()

        transposed_data.append(df)

    chartInfo['data'] = transposed_data

    temp = chartInfo['column_names_as_labels']
    chartInfo['column_names_as_labels'] = chartInfo['first_column_as_labels']
    chartInfo['first_column_as_labels'] = temp

    return chartInfo


def __get_dataframes(data):
    if (not isinstance(data, pd.DataFrame) and not __iterable(data)):
        return None

    if (isinstance(data, pd.DataFrame)):
        dfs = [data]
    else:
        for dataframe in data:
            if not isinstance(dataframe, pd.DataFrame):
                return None
        dfs = data

    return dfs


def __insert_object(slide, placeholder, chart):

    data = chart.get('data')

    if (data is None):
        return 'No data was supplied for chart'

    if (isinstance(data, pd.DataFrame)):
        chart['data'] = [data]

    for dataframe in chart['data']:
        if not isinstance(dataframe, pd.DataFrame):
            return 'Data supplied was neither a Pandas DataFrame, nor an array of Pandas DataFrames'

    if not isinstance(chart.get('column_names_as_labels'), bool):
        chart['column_names_as_labels'] = __infer_series_labels(
            chart['data'])

    if not isinstance(chart.get('first_column_as_labels'), bool):
        chart['first_column_as_labels'] = __infer_category_labels(
            chart['data'])

    transpose = chart.get('transpose', False)
    if transpose:
        chart = __transpose_data(chart)

    data = __get_dataframes(chart.get('data'))
    dataframe = data[0]

    chart_type = chart.get('chart_type', 'Table')

    if chart_type == CHART_TYPE.AREA.value:
        return __insert_chart(XL_CHART_TYPE.AREA, slide, placeholder, chart)
    elif chart_type == CHART_TYPE.AREA_STACKED.value:
        return __insert_chart(XL_CHART_TYPE.AREA_STACKED,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.AREA_STACKED_100.value:
        return __insert_chart(XL_CHART_TYPE.AREA_STACKED_100,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.BAR.value:
        return __insert_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.BAR_STACKED.value:
        return __insert_chart(XL_CHART_TYPE.BAR_STACKED,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.BAR_STACKED_100.value:
        return __insert_chart(XL_CHART_TYPE.BAR_STACKED_100,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.COLUMN.value:
        return __insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.COLUMN_STACKED.value:
        return __insert_chart(XL_CHART_TYPE.COLUMN_STACKED,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.COLUMN_STACKED_100.value:
        return __insert_chart(XL_CHART_TYPE.COLUMN_STACKED_100,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.LINE.value:
        return __insert_chart(XL_CHART_TYPE.LINE, slide, placeholder, chart)
    elif chart_type == CHART_TYPE.LINE_STACKED.value:
        return __insert_chart(XL_CHART_TYPE.LINE_STACKED,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.LINE_STACKED_100.value:
        return __insert_chart(XL_CHART_TYPE.LINE_STACKED_100,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.LINE_MARKED.value:
        return __insert_chart(XL_CHART_TYPE.LINE_MARKERS,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.LINE_MARKED_STACKED.value:
        return __insert_chart(XL_CHART_TYPE.LINE_MARKERS_STACKED,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.LINE_MARKED_STACKED_100.value:
        return __insert_chart(XL_CHART_TYPE.LINE_MARKERS_STACKED_100,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.DOUGHNUT.value:
        return __insert_chart(XL_CHART_TYPE.DOUGHNUT, slide, placeholder, chart)
    elif chart_type == CHART_TYPE.DOUGHNUT_EXPLODED.value:
        return __insert_chart(XL_CHART_TYPE.DOUGHNUT_EXPLODED,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.PIE.value:
        return __insert_chart(XL_CHART_TYPE.PIE, slide, placeholder, chart)
    elif chart_type == CHART_TYPE.PIE_EXPLODED.value:
        return __insert_chart(XL_CHART_TYPE.PIE_EXPLODED,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.RADAR.value:
        return __insert_chart(XL_CHART_TYPE.RADAR, slide, placeholder, chart)
    elif chart_type == CHART_TYPE.RADAR_FILLED.value:
        return __insert_chart(XL_CHART_TYPE.RADAR_FILLED,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.RADAR_MARKED.value:
        return __insert_chart(XL_CHART_TYPE.RADAR_MARKERS,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.XY_SCATTER.value:
        return __insert_xyzchart(XL_CHART_TYPE.XY_SCATTER,
                                 slide, placeholder, chart)
    elif chart_type == CHART_TYPE.XY_SCATTER_LINES.value:
        return __insert_xyzchart(XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
                                 slide, placeholder, chart)
    elif chart_type == CHART_TYPE.XY_SCATTER_LINES_SMOOTHED.value:
        return __insert_xyzchart(XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS,
                                 slide, placeholder, chart)
    elif chart_type == CHART_TYPE.XY_SCATTER_LINES_MARKED.value:
        return __insert_chart(XL_CHART_TYPE.XY_SCATTER_LINES,
                              slide, placeholder, chart)
    elif chart_type == CHART_TYPE.XY_SCATTER_LINES_MARKED_SMOOTHED.value:
        return __insert_xyzchart(XL_CHART_TYPE.XY_SCATTER_SMOOTH,
                                 slide, placeholder, chart)
    elif chart_type == CHART_TYPE.BUBBLE.value:
        return __insert_xyzchart(XL_CHART_TYPE.BUBBLE, slide, placeholder, chart)
    else:
        return __insert_table(slide, placeholder, chart)


def __insert_table(slide, placeholder, chartInfo):
    df = chartInfo['data'][0]

    columns = df.shape[1]
    rows = df.shape[0]
    if chartInfo['column_names_as_labels']:
        rows += 1

    # Create new element with same shape and position as placeholder
    table = slide.shapes.add_table(
        rows, columns, placeholder.left, placeholder.top, placeholder.width, placeholder.height).table
    table.first_row = chartInfo['column_names_as_labels']
    table.first_col = chartInfo['first_column_as_labels']

    # Populate table
    colNames = df.columns.tolist()

    rowNum = 0

    if chartInfo['column_names_as_labels']:
        col = 0
        for colName in colNames:
            table.cell(0, col).text = str(colName)
            col += 1
        rowNum += 1

    for index, row in df.iterrows():
        col = 0
        for colName in colNames:
            table.cell(rowNum, col).text = str(row.iloc[col])
            col += 1
        rowNum += 1

    return table


def __iterable(obj):
    return isinstance(obj, Iterable)


def __create_chartdata(chart):
    chart_data = CategoryChartData()

    # TODO: Deal with First Row as Labels and Column Names as Labels

    colNames = chart['data'][0].columns.tolist()
    offset = 0

    if (chart['first_column_as_labels']):
        offset = 1

    if len(colNames) > offset:

        colNum = 1
        for colName in colNames[offset:]:
            if (chart['column_names_as_labels']):
                chart_data.categories.add_category(colName)
            else:
                chart_data.categories.add_category('Category '+str(colNum))

        rowNum = 1
        for index, row in chart['data'][0].iterrows():
            data = []
            for colName in colNames[offset:]:
                data.append(row[colName])

            if chart['first_column_as_labels']:
                chart_data.add_series(str(row[0]), data)
            else:
                chart_data.add_series('Series ' + str(rowNum), data)

    return chart_data


def __create_xyzdata(dfs):
    chart_data = None

    seriesNum = 1

    for df in dfs:
        colNames = df.columns.tolist()
        name = 'Series ' + str(seriesNum)
        if hasattr(df, 'name') and df.name != "":
            name = df.name

        if len(colNames) > 1 and len(colNames) < 4:
            if len(colNames) == 2 and chart_data is None:
                chart_data = XyChartData()
            elif len(colNames) == 3 and chart_data is None:
                chart_data = BubbleChartData()

            series = chart_data.add_series(name)
            for index, row in df.iterrows():
                data = []
                for colName in colNames:
                    data.append(row[colName])

                if len(colNames) == 2:
                    series.add_data_point(data[0], data[1])
                else:
                    series.add_data_point(data[0], data[1], data[2])

            seriesNum += 1

    return chart_data


def __insert_chart(chart_type, slide, placeholder, chart):
    chart_data = __create_chartdata(chart)
    if chart_data is None:
        return 'Could not create chart data'

    # Create new element with same shape and position as placeholder
    new_chart = slide.shapes.add_chart(chart_type, placeholder.left,
                                       placeholder.top, placeholder.width, placeholder.height, chart_data).chart

    __set_font_object(new_chart.font, chart.get('body_font'))

    __set_chart_title(new_chart, chart)

    __set_axis_object(new_chart.value_axis, chart.get('y_axis'))

    __set_chart_legend(new_chart, chart)

    return new_chart


def __set_chart_title(new_chart, chart):
    title = chart.get('title')
    if title is not None:
        title_tf = new_chart.chart_title.text_frame
        title_tf.clear()
        title_p = title_tf.paragraphs[0]
        title_p.add_run().text = title


def __set_chart_legend(new_chart, chart):
    legend_position = chart.get('legend_position')
    if legend_position is not None and legend_position != LEGEND_POSITION.NONE.value:
        new_chart.has_legend = True
        if legend_position == LEGEND_POSITION.BOTTOM.value:
            new_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        elif legend_position == LEGEND_POSITION.CORNER.value:
            new_chart.legend.position = XL_LEGEND_POSITION.CORNER
        elif legend_position == LEGEND_POSITION.LEFT.value:
            new_chart.legend.position = XL_LEGEND_POSITION.LEFT
        elif legend_position == LEGEND_POSITION.RIGHT.value:
            new_chart.legend.position = XL_LEGEND_POSITION.RIGHT
        elif legend_position == LEGEND_POSITION.TOP.value:
            new_chart.legend.position = XL_LEGEND_POSITION.TOP

        if chart.get('overlay_legend', False):
            new_chart.legend.include_in_layout = True
        else:
            new_chart.legend.include_in_layout = False


def __set_font_object(font_object, font):
    font_object.name = font['name']
    font_object.size = Pt(font['size'])


def __set_axis_object(axis_object, axis):
    if axis is None:
        axis = dict()

    axis_object.visible = axis.get('visible', True) == True
    axis_object.minimum_scale = axis.get('minimum_scale')
    axis_object.maximum_scale = axis.get('maximum_scale')
    has_major_gridlines = axis.get('has_major_grid_lines', False)
    axis_object.has_major_gridlines = has_major_gridlines
    has_minor_gridlines = axis.get('has_minor_grid_lines', False)
    axis_object.has_minor_gridlines = has_minor_gridlines
    has_title = axis.get('title', False) != False
    if has_title:
        axis_object.has_title = True
        axis_object.axis_title = axis.get('title')

    axis_object.tick_labels.number_format = axis.get(
        'number_format', '$#0.0,,"M";[Red]($#0.0,,"M")')


def __insert_xyzchart(chart_type, slide, placeholder, chart):
    chart_data = __create_xyzdata(chart['data'])
    if chart_data is None:
        return 'Could not create chart data'

    # Create new element with same shape and position as placeholder
    new_chart = slide.shapes.add_chart(chart_type, placeholder.left,
                                       placeholder.top, placeholder.width, placeholder.height, chart_data).chart

    __set_font_object(new_chart.font, chart.get('body_font'))

    __set_chart_title(new_chart, chart)

    __set_axis_object(new_chart.value_axis, chart.get('x_axis'))

    __set_axis_object(new_chart.value_axis, chart.get('y_axis'))

    __set_chart_legend(new_chart, chart)

    return new_chart


def __get_datafile_name(filename):
    """
    return the default template file that comes with the package
    """
    return Path(__file__).parent / "data/" + filename
